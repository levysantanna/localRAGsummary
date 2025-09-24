"""
Document processor for handling various file types with OCR and multilingual support
"""
import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
import mimetypes
import hashlib
from datetime import datetime

# Document processing imports
import PyPDF2 as pypdf2
import pdfplumber
from PIL import Image
import pytesseract
import easyocr
import cv2
import numpy as np
from docx import Document
from pptx import Presentation

# ODF processor
from processors.odf_processor import ODFProcessor

# Code analysis imports
import tree_sitter
from tree_sitter import Language, Parser

# Text processing
import nltk
import spacy
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import re

from config import *

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Main document processor class"""
    
    def __init__(self):
        self.setup_nltk()
        self.setup_spacy()
        self.setup_ocr()
        self.setup_code_parsers()
        self.odf_processor = ODFProcessor()
        
    def setup_nltk(self):
        """Setup NLTK resources"""
        try:
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
            nltk.download('portuguese', quiet=True)
        except Exception as e:
            logger.warning(f"NLTK setup warning: {e}")
    
    def setup_spacy(self):
        """Setup spaCy models"""
        try:
            # Try to load Portuguese model
            self.nlp_pt = spacy.load("pt_core_news_sm")
        except OSError:
            logger.warning("Portuguese spaCy model not found. Install with: python -m spacy download pt_core_news_sm")
            self.nlp_pt = None
    
    def setup_ocr(self):
        """Setup OCR engines"""
        self.easyocr_reader = easyocr.Reader(['pt', 'en'], gpu=DEVICE_CONFIG['use_gpu'])
        
    def setup_code_parsers(self):
        """Setup code parsers for different languages"""
        self.code_parsers = {}
        try:
            # Python parser
            PY_LANGUAGE = Language('tree_sitter_python.so', 'python')
            self.code_parsers['python'] = Parser(PY_LANGUAGE)
            
            # JavaScript parser
            JS_LANGUAGE = Language('tree_sitter_javascript.so', 'javascript')
            self.code_parsers['javascript'] = Parser(JS_LANGUAGE)
            
            # Java parser
            JAVA_LANGUAGE = Language('tree_sitter_java.so', 'java')
            self.code_parsers['java'] = Parser(JAVA_LANGUAGE)
            
        except Exception as e:
            logger.warning(f"Code parser setup warning: {e}")
    
    def process_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a single document and extract content"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file metadata
        metadata = self._get_file_metadata(file_path)
        
        # Determine file type and process accordingly
        file_type = self._get_file_type(file_path)
        
        try:
            if file_type == 'pdf':
                content = self._process_pdf(file_path)
            elif file_type == 'image':
                content = self._process_image(file_path)
            elif file_type == 'document':
                content = self._process_document_file(file_path)
            elif file_type == 'odf':
                content = self._process_odf_file(file_path)
            elif file_type == 'code':
                content = self._process_code_file(file_path)
            elif file_type == 'text':
                content = self._process_text_file(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_path}")
                return None
            
            # Post-process content
            processed_content = self._post_process_content(content, metadata)
            
            return {
                'file_path': str(file_path),
                'file_type': file_type,
                'metadata': metadata,
                'content': processed_content,
                'processed_at': datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return None
    
    def _get_file_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract file metadata"""
        stat = file_path.stat()
        return {
            'filename': file_path.name,
            'size_bytes': stat.st_size,
            'size_mb': round(stat.st_size / (1024 * 1024), 2),
            'created_at': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'modified_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'file_hash': self._calculate_file_hash(file_path)
        }
    
    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate file hash for deduplication"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def _get_file_type(self, file_path: Path) -> str:
        """Determine file type based on extension"""
        suffix = file_path.suffix.lower()
        
        for file_type, extensions in SUPPORTED_EXTENSIONS.items():
            if suffix in extensions:
                return file_type
        
        return 'unknown'
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF files with OCR for scanned documents"""
        content = {
            'text': '',
            'images': [],
            'tables': [],
            'metadata': {}
        }
        
        try:
            # Try text extraction first
            with pdfplumber.open(file_path) as pdf:
                text_content = []
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        for table_num, table in enumerate(tables):
                            content['tables'].append({
                                'page': page_num + 1,
                                'table': table_num + 1,
                                'data': table
                            })
                
                content['text'] = '\n\n'.join(text_content)
                
                # If no text found, try OCR
                if not content['text'].strip():
                    logger.info(f"No text found in PDF, attempting OCR: {file_path}")
                    content = self._ocr_pdf_pages(file_path)
        
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            # Fallback to OCR
            content = self._ocr_pdf_pages(file_path)
        
        return content
    
    def _ocr_pdf_pages(self, file_path: Path) -> Dict[str, Any]:
        """Perform OCR on PDF pages"""
        content = {'text': '', 'images': []}
        
        try:
            import fitz  # PyMuPDF
            import io
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                
                # Convert to PIL Image
                img = Image.open(io.BytesIO(img_data))
                
                # Perform OCR
                ocr_result = self._perform_ocr(img)
                if ocr_result:
                    content['text'] += f"\n--- Page {page_num + 1} ---\n{ocr_result}\n"
            
            doc.close()
            
        except ImportError:
            logger.warning("PyMuPDF not available for PDF OCR")
        except Exception as e:
            logger.error(f"OCR error for PDF {file_path}: {e}")
        
        return content
    
    def _process_image(self, file_path: Path) -> Dict[str, Any]:
        """Process image files with OCR"""
        content = {'text': '', 'metadata': {}}
        
        try:
            # Load image
            img = Image.open(file_path)
            content['metadata'] = {
                'format': img.format,
                'mode': img.mode,
                'size': img.size
            }
            
            # Perform OCR
            ocr_result = self._perform_ocr(img)
            if ocr_result:
                content['text'] = ocr_result
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
        
        return content
    
    def _perform_ocr(self, img: Image.Image) -> str:
        """Perform OCR on image using multiple engines"""
        try:
            # Try EasyOCR first (better for Portuguese)
            results = self.easyocr_reader.readtext(np.array(img))
            text_parts = []
            
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Filter low confidence results
                    text_parts.append(text)
            
            if text_parts:
                return ' '.join(text_parts)
            
            # Fallback to Tesseract
            return pytesseract.image_to_string(img, lang='por+eng')
            
        except Exception as e:
            logger.error(f"OCR error: {e}")
            return ""
    
    def _process_document_file(self, file_path: Path) -> Dict[str, Any]:
        """Process Word and PowerPoint files"""
        content = {'text': '', 'metadata': {}}
        
        try:
            if file_path.suffix.lower() in ['.docx']:
                doc = Document(file_path)
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                content['text'] = '\n'.join(paragraphs)
                
            elif file_path.suffix.lower() in ['.pptx']:
                prs = Presentation(file_path)
                slides = []
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        slides.append(f"Slide {slide_num + 1}: {' '.join(slide_text)}")
                content['text'] = '\n'.join(slides)
        
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
        
        return content
    
    def _process_odf_file(self, file_path: Path) -> Dict[str, Any]:
        """Process Open Document Format files (ODT, ODS, ODP)"""
        try:
            if self.odf_processor.can_process(str(file_path)):
                result = self.odf_processor.process_odf_document(str(file_path))
                return {
                    'text': result['content']['text'],
                    'metadata': result['metadata'],
                    'processing_info': result['processing_info']
                }
            else:
                logger.warning(f"ODF processor cannot handle {file_path}")
                return {'text': '', 'metadata': {}, 'processing_info': {'error': 'Unsupported ODF format'}}
        
        except Exception as e:
            logger.error(f"Error processing ODF file {file_path}: {e}")
            return {'text': '', 'metadata': {}, 'processing_info': {'error': str(e)}}
    
    def _process_code_file(self, file_path: Path) -> Dict[str, Any]:
        """Process code files with syntax analysis"""
        content = {'text': '', 'code_structure': {}, 'metadata': {}}
        
        try:
            # Read file content
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                code_content = f.read()
            
            content['text'] = code_content
            
            # Analyze code structure if parser available
            language = self._get_code_language(file_path)
            if language in self.code_parsers:
                structure = self._analyze_code_structure(code_content, language)
                content['code_structure'] = structure
        
        except Exception as e:
            logger.error(f"Error processing code file {file_path}: {e}")
        
        return content
    
    def _get_code_language(self, file_path: Path) -> str:
        """Determine programming language from file extension"""
        extension_map = {
            '.py': 'python',
            '.js': 'javascript',
            '.ts': 'typescript',
            '.java': 'java',
            '.cpp': 'cpp',
            '.c': 'c',
            '.cs': 'csharp',
            '.php': 'php',
            '.rb': 'ruby',
            '.go': 'go',
            '.rs': 'rust',
            '.swift': 'swift',
            '.kt': 'kotlin',
            '.scala': 'scala',
            '.r': 'r',
            '.m': 'matlab',
            '.pl': 'perl',
            '.sh': 'bash',
            '.sql': 'sql',
            '.html': 'html',
            '.css': 'css',
            '.xml': 'xml',
            '.json': 'json',
            '.yaml': 'yaml',
            '.yml': 'yaml'
        }
        
        return extension_map.get(file_path.suffix.lower(), 'unknown')
    
    def _analyze_code_structure(self, code: str, language: str) -> Dict[str, Any]:
        """Analyze code structure using tree-sitter"""
        try:
            parser = self.code_parsers[language]
            tree = parser.parse(bytes(code, 'utf8'))
            
            # Extract functions, classes, imports, etc.
            structure = {
                'functions': [],
                'classes': [],
                'imports': [],
                'comments': []
            }
            
            # This is a simplified version - would need language-specific implementation
            # For now, return basic structure
            return structure
            
        except Exception as e:
            logger.error(f"Code analysis error: {e}")
            return {}
    
    def _process_text_file(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text files"""
        content = {'text': '', 'metadata': {}}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content['text'] = f.read()
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
        
        return content
    
    def _post_process_content(self, content: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Post-process extracted content"""
        if not content.get('text'):
            return content
        
        # Clean and normalize text
        text = content['text']
        
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Language detection and processing
        if self.nlp_pt:
            doc = self.nlp_pt(text[:1000])  # Sample for language detection
            # Add language-specific processing here
        
        content['text'] = text.strip()
        
        # Add word count and other statistics
        content['statistics'] = {
            'word_count': len(text.split()),
            'character_count': len(text),
            'language': 'pt'  # Default to Portuguese
        }
        
        return content
    
    def process_directory(self, directory_path: Union[str, Path], recursive: bool = True) -> List[Dict[str, Any]]:
        """Process all documents in a directory"""
        directory_path = Path(directory_path)
        processed_documents = []
        
        if not directory_path.exists():
            logger.error(f"Directory not found: {directory_path}")
            return processed_documents
        
        # Get all files to process
        files_to_process = []
        if recursive:
            for ext in [ext for exts in SUPPORTED_EXTENSIONS.values() for ext in exts]:
                files_to_process.extend(directory_path.rglob(f"*{ext}"))
        else:
            for ext in [ext for exts in SUPPORTED_EXTENSIONS.values() for ext in exts]:
                files_to_process.extend(directory_path.glob(f"*{ext}"))
        
        logger.info(f"Found {len(files_to_process)} files to process")
        
        # Process files
        for file_path in files_to_process:
            try:
                if file_path.stat().st_size > PROCESSING_CONFIG['max_file_size_mb'] * 1024 * 1024:
                    logger.warning(f"File too large, skipping: {file_path}")
                    continue
                
                logger.info(f"Processing: {file_path}")
                result = self.process_document(file_path)
                
                if result:
                    processed_documents.append(result)
                
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")
        
        logger.info(f"Successfully processed {len(processed_documents)} documents")
        return processed_documents
